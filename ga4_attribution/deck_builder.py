"""Build a PowerPoint deck from GA4 attribution results."""
from __future__ import annotations

import io
from datetime import datetime

import pandas as pd
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

# ── Brand palette ────────────────────────────────────────────────────────────
_BLUE  = RGBColor(0x45, 0x7B, 0x9D)
_DARK  = RGBColor(0x1A, 0x1A, 0x2E)
_GRAY  = RGBColor(0xF0, 0xF4, 0xF8)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_MID   = RGBColor(0x8A, 0x9B, 0xB0)
_CODE_BG  = RGBColor(0x1E, 0x1E, 0x2E)
_CODE_FG  = RGBColor(0xA8, 0xD8, 0xA8)

MODEL_COLORS: dict[str, RGBColor] = {
    "last_touch":     RGBColor(0xE6, 0x39, 0x46),
    "first_touch":    RGBColor(0x45, 0x7B, 0x9D),
    "linear":         RGBColor(0x2A, 0x9D, 0x8F),
    "time_decay":     RGBColor(0xE9, 0xC4, 0x6A),
    "position_based": RGBColor(0x9B, 0x5D, 0xE5),
    "shapley":        RGBColor(0xF4, 0xA2, 0x61),
    "markov":         RGBColor(0x06, 0xD6, 0xA0),
}

MODEL_DESCRIPTIONS = {
    "last_touch":     "100% credit to the final touchpoint before conversion",
    "first_touch":    "100% credit to the first touchpoint in the journey",
    "linear":         "Equal credit split across all touchpoints",
    "time_decay":     "More credit to touchpoints closer to conversion",
    "position_based": "40% first / 40% last / 20% middle (U-shape)",
    "shapley":        "Game-theory marginal contribution per channel",
    "markov":         "Data-driven removal effects via transition matrix",
}

W = Inches(13.33)
H = Inches(7.5)
MARGIN = Inches(0.6)

# ── Helpers ──────────────────────────────────────────────────────────────────

def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # fully blank


def _rect(slide, left, top, width, height, fill_color: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _text(
    slide,
    left, top, width, height,
    text: str,
    size: int = 12,
    bold: bool = False,
    color: RGBColor = _DARK,
    wrap: bool = True,
    font: str | None = None,
):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if font:
        run.font.name = font
    return txb


def _header_bar(slide, title: str, dark: bool = False):
    bg = _DARK if dark else _BLUE
    _rect(slide, 0, 0, W, Inches(1.05), bg)
    _text(slide, MARGIN, Inches(0.18), W - MARGIN, Inches(0.7),
          title, size=22, bold=True, color=_WHITE)

# ── Slides ───────────────────────────────────────────────────────────────────

def _title_slide(prs: Presentation, meta: dict) -> None:
    slide = _blank_slide(prs)

    # Left accent bar
    _rect(slide, 0, 0, Inches(0.35), H, _BLUE)

    _text(slide, Inches(0.85), Inches(2.0), Inches(11.5), Inches(1.4),
          "GA4 Multi-Touch Attribution Analysis",
          size=34, bold=True, color=_DARK)

    subtitle = (
        f"Period: {meta.get('start_date', '')}  →  {meta.get('end_date', '')}\n"
        f"Conversions: {meta.get('total_conversions', 0):,}   ·   "
        f"Total value: ${meta.get('total_conversion_value', 0):,.2f}   ·   "
        f"Avg path: {meta.get('avg_path_length', 0):.1f} touchpoints"
    )
    _text(slide, Inches(0.85), Inches(3.7), Inches(11.5), Inches(1.0),
          subtitle, size=14, color=_MID)

    _text(slide, Inches(0.85), Inches(6.9), Inches(9), Inches(0.4),
          f"Generated {datetime.today().strftime('%B %d, %Y')}  ·  GA4 Attribution Agent",
          size=9, color=_MID)


def _metrics_slide(prs: Presentation, meta: dict, claude_summary: str) -> None:
    slide = _blank_slide(prs)
    _header_bar(slide, "Key Metrics & Summary")

    cards = [
        ("Conversions",    f"{meta.get('total_conversions', 0):,}"),
        ("Total Value",    f"${meta.get('total_conversion_value', 0):,.2f}"),
        ("Avg Path Length",f"{meta.get('avg_path_length', 0):.1f} touches"),
    ]
    card_w = Inches(3.7)
    card_h = Inches(1.45)
    card_top = Inches(1.2)
    lefts = [MARGIN, Inches(4.82), Inches(9.03)]

    for (label, value), left in zip(cards, lefts):
        _rect(slide, left, card_top, card_w, card_h, _GRAY)
        _text(slide, left + Inches(0.18), card_top + Inches(0.12),
              card_w - Inches(0.3), Inches(0.42), label, size=10, color=_MID)
        _text(slide, left + Inches(0.18), card_top + Inches(0.55),
              card_w - Inches(0.3), Inches(0.75), value, size=22, bold=True, color=_DARK)

    if claude_summary:
        _text(slide, MARGIN, Inches(2.85), W - MARGIN * 2, Inches(0.38),
              "Analysis Summary", size=12, bold=True, color=_DARK)
        summary = claude_summary[:750] + ("…" if len(claude_summary) > 750 else "")
        _text(slide, MARGIN, Inches(3.25), W - MARGIN * 2, Inches(3.9),
              summary, size=10, color=_DARK)


def _table_slide(prs: Presentation, results_df: pd.DataFrame) -> None:
    slide = _blank_slide(prs)
    _header_bar(slide, "Attribution Results by Channel")

    cols = list(results_df.columns)
    n_rows = len(results_df) + 1
    n_cols = len(cols)

    tbl_left  = MARGIN
    tbl_top   = Inches(1.2)
    tbl_w     = W - MARGIN * 2
    row_h     = Inches(0.42)
    tbl_h     = row_h * n_rows

    table = slide.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_w, tbl_h).table

    # Column widths: channel col wider, model cols equal
    ch_w = Inches(2.0)
    model_w = int((tbl_w - ch_w) / (n_cols - 1))
    table.columns[0].width = int(ch_w)
    for i in range(1, n_cols):
        table.columns[i].width = model_w

    # Header
    for j, col in enumerate(cols):
        cell = table.cell(0, j)
        cell.text = col.replace("_", " ").title()
        cell.fill.solid()
        cell.fill.fore_color.rgb = _BLUE
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True
                run.font.size = Pt(9)
                run.font.color.rgb = _WHITE

    # Data rows
    for i, (_, row) in enumerate(results_df.iterrows()):
        bg = _GRAY if i % 2 == 0 else _WHITE
        for j, col in enumerate(cols):
            cell = table.cell(i + 1, j)
            val = row[col]
            cell.text = str(val) if j == 0 else f"{val:,.1f}"
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(9)
                    run.font.color.rgb = _DARK


def _chart_slide(prs: Presentation, results_df: pd.DataFrame, meta: dict) -> None:
    slide = _blank_slide(prs)
    _header_bar(slide, "Model Comparison by Channel")

    model_cols = [c for c in results_df.columns if c != "channel"]
    channels = results_df["channel"].tolist()

    chart_data = ChartData()
    chart_data.categories = channels
    for model in model_cols:
        chart_data.add_series(
            model.replace("_", " ").title(),
            [round(v, 2) for v in results_df[model].tolist()],
        )

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        MARGIN, Inches(1.15),
        W - MARGIN * 2, Inches(6.1),
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.position = 2  # bottom
    chart.legend.include_in_layout = False

    for i, model in enumerate(model_cols):
        series = chart.series[i]
        color = MODEL_COLORS.get(model, RGBColor(0x88, 0x88, 0x88))
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color


def _model_guide_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _header_bar(slide, "Attribution Model Guide")

    items = list(MODEL_DESCRIPTIONS.items())
    col_w = Inches(5.9)

    for idx, (model, desc) in enumerate(items):
        col = idx % 2
        row = idx // 2
        left = MARGIN + col * Inches(6.75)
        top  = Inches(1.25) + row * Inches(1.3)

        color = MODEL_COLORS.get(model, _BLUE)
        _rect(slide, left, top, Inches(0.1), Inches(0.85), color)
        _text(slide, left + Inches(0.22), top + Inches(0.05),
              col_w, Inches(0.4),
              model.replace("_", " ").title(), size=12, bold=True, color=_DARK)
        _text(slide, left + Inches(0.22), top + Inches(0.44),
              col_w, Inches(0.45), desc, size=10, color=_MID)


def _sql_slide(prs: Presentation, sql: str) -> None:
    slide = _blank_slide(prs)
    _header_bar(slide, "Generated SQL — Journey Extraction Query", dark=True)

    _rect(slide, MARGIN, Inches(1.18), W - MARGIN * 2, Inches(6.08), _CODE_BG)

    preview = sql[:1400] + ("\n\n-- [truncated]" if len(sql) > 1400 else "")
    _text(slide,
          MARGIN + Inches(0.15), Inches(1.28),
          W - MARGIN * 2 - Inches(0.3), Inches(5.9),
          preview, size=7.5, color=_CODE_FG, font="Courier New")

# ── Public API ───────────────────────────────────────────────────────────────

def build_deck(
    results_df: pd.DataFrame,
    meta: dict,
    sql: str,
    claude_summary: str = "",
) -> io.BytesIO:
    """Build a 6-slide PowerPoint deck and return as an in-memory BytesIO."""
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    _title_slide(prs, meta)
    _metrics_slide(prs, meta, claude_summary)
    _table_slide(prs, results_df)
    _chart_slide(prs, results_df, meta)
    _model_guide_slide(prs)
    _sql_slide(prs, sql)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
